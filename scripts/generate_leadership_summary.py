#!/usr/bin/env python3
"""
Generate a 2-slide PPTX summarizing the CI Agent project for leadership.
Slide 1: Current capabilities
Slide 2: Roadmap & timeline
Uses the same Lilly brand system as generate_deck.py / generate_architecture_slide.py.
"""

import os
from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ═══════════════════════════════════════════════════════════════════
# DESIGN CONSTANTS (Lilly brand)
# ═══════════════════════════════════════════════════════════════════

PRIMARY_COLOR = RGBColor(0xE1, 0x25, 0x1B)       # Lilly Red
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)           # Navy for boxes
MID_BLUE = RGBColor(0x2C, 0x3E, 0x5F)            # Slightly lighter navy
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)            # Light gray background
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x21, 0x21, 0x21)
GRAY_TEXT = RGBColor(0x66, 0x66, 0x66)
ACCENT_TEAL = RGBColor(0x00, 0x7B, 0x8A)         # Teal accent
DONE_GREEN = RGBColor(0x2E, 0x7D, 0x32)          # Green for completed phase

HEADER_FONT = "Arial"
BODY_FONT = "Arial"

SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.50)


# ═══════════════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════════════

def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=11, font_color=WHITE, bold=False,
                     align=PP_ALIGN.LEFT, border_color=None, v_anchor=MSO_ANCHOR.MIDDLE):
    """Add a rounded rectangle with optional multi-line text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = v_anchor
        tf.margin_left = Pt(10)
        tf.margin_right = Pt(10)
        tf.margin_top = Pt(6)
        tf.margin_bottom = Pt(6)
        lines = text.split("\n")
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(font_size) if i > 0 else Pt(font_size + 1)
            p.font.name = BODY_FONT
            p.font.color.rgb = font_color
            p.font.bold = bold if i == 0 else False
            p.alignment = align
            p.space_after = Pt(2)
    return shape


def add_text(slide, left, top, width, height, text, font_size=11,
             color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT):
    """Add a plain text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = BODY_FONT
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox


def add_multiline_text(slide, left, top, width, height, lines_data,
                       default_size=10, default_color=DARK_TEXT, align=PP_ALIGN.LEFT):
    """Add a text box with multiple lines, each with optional formatting.
    lines_data: list of tuples (text, font_size, color, bold)
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold) in enumerate(lines_data):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size if size else default_size)
        p.font.name = BODY_FONT
        p.font.color.rgb = color if color else default_color
        p.font.bold = bold
        p.alignment = align
        p.space_after = Pt(4)
    return txBox


def add_arrow_right(slide, left, top, width, height):
    """Add a right-pointing arrow."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GRAY_TEXT
    shape.line.fill.background()
    return shape


def add_title_banner(slide, title_text):
    """Add the red title banner at the top of a slide."""
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.9)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = PRIMARY_COLOR
    banner.line.fill.background()

    add_text(slide, Inches(0.5), Inches(0.1), Inches(12), Inches(0.8),
             title_text, font_size=26, color=WHITE, bold=True)


def add_footer(slide, text):
    """Add a footer line at the bottom."""
    add_text(slide, Inches(0.4), Inches(7.1), Inches(12), Inches(0.3),
             text, font_size=8, color=GRAY_TEXT)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1: CURRENT CAPABILITIES
# ═══════════════════════════════════════════════════════════════════

def build_slide_capabilities(prs):
    """Build the 'Current Capabilities' slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    add_title_banner(slide, "CI Agent — Current Capabilities")

    # Subtitle
    y = 1.05
    add_text(slide, Inches(0.4), Inches(y), Inches(12), Inches(0.4),
             "An AI-powered agent that extracts, structures, analyzes, and summarizes "
             "competitor clinical trial data for Immunology.",
             font_size=12, color=GRAY_TEXT)

    # ─── 5 Capability Boxes ───
    y = 1.55
    add_text(slide, Inches(0.4), Inches(y), Inches(6), Inches(0.3),
             "CORE CAPABILITIES", font_size=10, color=GRAY_TEXT, bold=True)

    y += 0.35
    box_w = Inches(2.35)
    box_h = Inches(1.45)
    gap = Inches(0.18)
    x_start = Inches(0.4)

    capabilities = [
        ("Data Extraction",
         "Press releases, CTgov,\nPubMed, pasted text\n→ structured table\n(arms × endpoints)"),
        ("BNMA Ridge Plot",
         "Batman NMA output\n→ ridge/density plot\nwith compound\nrecommendation engine"),
        ("Slide Deck Gen",
         "Quick (5-slide) or\nDetailed (8+N slide)\nwith figures, BNMA,\nspeaker notes"),
        ("Pipeline Mode",
         "End-to-end:\nextract → BNMA →\nlandscape → slides\n→ QA verification"),
        ("Competitor DB",
         "SQLite + REST API\ndeployed on Posit\nConnect; queryable,\ndownloadable (Excel)"),
    ]

    for i, (title, desc) in enumerate(capabilities):
        x = x_start + i * (box_w + gap)
        add_rounded_rect(slide, x, Inches(y), box_w, box_h,
                         DARK_BLUE, f"{title}\n{desc}",
                         font_size=9, font_color=WHITE, bold=True,
                         align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.TOP)

    # ─── "Demonstrated On" callout ───
    y += 1.65
    add_text(slide, Inches(0.4), Inches(y), Inches(6), Inches(0.3),
             "DEMONSTRATED ON", font_size=10, color=GRAY_TEXT, bold=True)

    y += 0.30
    add_rounded_rect(slide, Inches(0.4), Inches(y), Inches(12.5), Inches(0.65),
                     LIGHT_BG,
                     "Zumilokibart (APG777) APEX Phase 2B — Atopic Dermatitis  •  "
                     "Generated: detailed presenter deck (10 slides) + BNMA ridge plots (EASI-75, IGA 0/1)",
                     font_size=10, font_color=DARK_TEXT, bold=False,
                     align=PP_ALIGN.LEFT, border_color=ACCENT_TEAL)

    # ─── Technical Details Row ───
    y += 0.90
    add_text(slide, Inches(0.4), Inches(y), Inches(6), Inches(0.3),
             "TECHNICAL DETAILS", font_size=10, color=GRAY_TEXT, bold=True)

    y += 0.30
    details = [
        ("Platform:", "Claude Code (VS Code extension or terminal)"),
        ("Languages:", "Python (slide gen) + R (BNMA ridge plot) + SQLite/Plumber (DB API)"),
        ("Indications:", "AD, Psoriasis, UC, RA, CRSwNP, PsA, Crohn's, SLE, Asthma, COPD, IPF, AR"),
        ("Config-driven:", "JSON configs per compound — reusable, no code edits needed per analysis"),
    ]

    for i, (label, value) in enumerate(details):
        x_label = Inches(0.5)
        x_value = Inches(2.4)
        row_y = Inches(y + i * 0.28)
        add_text(slide, x_label, row_y, Inches(2.0), Inches(0.25),
                 label, font_size=9, color=DARK_TEXT, bold=True)
        add_text(slide, x_value, row_y, Inches(10.5), Inches(0.25),
                 value, font_size=9, color=DARK_TEXT, bold=False)

    # Footer
    add_footer(slide,
               "Lilly Immunology Competitive Intelligence  •  Built Jul–Aug 2026  •  "
               "Review is required before disclosure.")

    return slide


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2: ROADMAP & TIMELINE
# ═══════════════════════════════════════════════════════════════════

def build_slide_roadmap(prs):
    """Build the 'Roadmap & Timeline' slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    add_title_banner(slide, "CI Agent — Roadmap & Timeline")

    # ─── Timeline phases ───
    y = 1.15
    add_text(slide, Inches(0.4), Inches(y), Inches(10), Inches(0.3),
             "3-PHASE PLAN  •  Target completion: October 2026",
             font_size=11, color=GRAY_TEXT, bold=True)

    y += 0.45

    # Phase definitions
    phases = [
        {
            "title": "Phase 1 — Foundation  ✓",
            "timeframe": "Jul – Aug 2026",
            "color": DONE_GREEN,
            "status": "COMPLETE",
            "items": [
                "Data extraction from 5 source types",
                "BNMA ridge plot from Batman NMA output",
                "Slide deck generation (Quick + Detailed)",
                "Pipeline mode (end-to-end orchestration)",
                "Competitor DB + REST API on Posit Connect",
                "JSON config system (reusable per compound)",
            ],
        },
        {
            "title": "Phase 2 — Validate & Harden",
            "timeframe": "September 2026",
            "color": ACCENT_TEAL,
            "status": "NEXT",
            "items": [
                "Test across all 12 indications (beyond AD)",
                "Standalone landscape chart script (R/ggplot2)",
                "Team onboarding & training materials",
                "Robust error handling & retry logic",
                "Populate competitor DB with historical data",
            ],
        },
        {
            "title": "Phase 3 — Integrate & Deploy",
            "timeframe": "October 2026",
            "color": MID_BLUE,
            "status": "PLANNED",
            "items": [
                "CILand auto-fetch (MCP connector for SSO)",
                "New readout monitoring / alerts",
                "Team-facing web dashboard (Shiny/HTML)",
                "Multi-user access & permissions",
                "Full production deployment",
            ],
        },
    ]

    phase_w = Inches(3.95)
    phase_h = Inches(3.60)
    phase_gap = Inches(0.25)
    x_start = Inches(0.4)

    for i, phase in enumerate(phases):
        x = x_start + i * (phase_w + phase_gap)

        # Phase box (outer)
        add_rounded_rect(slide, x, Inches(y), phase_w, phase_h,
                         WHITE, "", border_color=phase["color"])

        # Phase header bar
        header_h = Inches(0.70)
        header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(y), phase_w, header_h
        )
        header.fill.solid()
        header.fill.fore_color.rgb = phase["color"]
        header.line.fill.background()

        # Header text
        tf = header.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Pt(12)
        p = tf.paragraphs[0]
        p.text = phase["title"]
        p.font.size = Pt(12)
        p.font.name = BODY_FONT
        p.font.color.rgb = WHITE
        p.font.bold = True
        p2 = tf.add_paragraph()
        p2.text = phase["timeframe"]
        p2.font.size = Pt(9)
        p2.font.name = BODY_FONT
        p2.font.color.rgb = WHITE
        p2.font.bold = False

        # Bullet items
        item_y = y + 0.80
        for j, item in enumerate(phase["items"]):
            add_text(slide, x + Inches(0.15), Inches(item_y + j * 0.38),
                     phase_w - Inches(0.3), Inches(0.35),
                     f"• {item}", font_size=9.5, color=DARK_TEXT)

        # Arrow between phases
        if i < len(phases) - 1:
            arrow_x = x + phase_w + Inches(0.02)
            arrow_y = Inches(y) + phase_h / 2 - Inches(0.1)
            add_arrow_right(slide, arrow_x, arrow_y, Inches(0.22), Inches(0.2))

    # ─── Value proposition at bottom ───
    y_bottom = 5.55
    add_text(slide, Inches(0.4), Inches(y_bottom), Inches(6), Inches(0.3),
             "VALUE", font_size=10, color=GRAY_TEXT, bold=True)

    y_bottom += 0.30
    value_items = [
        "• Reduces competitor readout turnaround from days → hours (extraction + deck generation automated)",
        "• Standardized BNMA-based indirect comparisons (avoids cross-trial naïve comparison pitfalls)",
        "• Reusable JSON configs — subsequent readouts for the same compound require no re-setup",
        "• Shared database enables team-wide access to structured competitor efficacy data",
    ]

    add_rounded_rect(slide, Inches(0.4), Inches(y_bottom), Inches(12.5), Inches(1.30),
                     LIGHT_BG, "\n".join(value_items),
                     font_size=9.5, font_color=DARK_TEXT, bold=False,
                     align=PP_ALIGN.LEFT, border_color=GRAY_TEXT,
                     v_anchor=MSO_ANCHOR.TOP)

    # Footer
    add_footer(slide,
               "Lilly Immunology Competitive Intelligence  •  "
               "Contact: [your name]  •  Live demo available upon request.")

    return slide


# ═══════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    build_slide_capabilities(prs)
    build_slide_roadmap(prs)

    # Save
    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    today = date.today().isoformat()
    output_path = os.path.join(base_dir, f"CI_Agent_Leadership_Summary_{today}.pptx")
    prs.save(output_path)
    print(f"✓ Saved: {output_path}")
    print(f"  Slides: 2")


if __name__ == "__main__":
    main()
