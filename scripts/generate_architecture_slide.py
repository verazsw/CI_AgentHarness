#!/usr/bin/env python3
"""
Generate a single-slide PPTX summarizing the CI Agent architecture/workflow.
Uses the same Lilly brand system as generate_deck.py.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ═══════════════════════════════════════════════════════════════════
# DESIGN CONSTANTS
# ═══════════════════════════════════════════════════════════════════

PRIMARY_COLOR = RGBColor(0xE1, 0x25, 0x1B)       # Lilly Red
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)           # Navy for boxes
MID_BLUE = RGBColor(0x2C, 0x3E, 0x5F)            # Slightly lighter navy
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)            # Light gray background
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x21, 0x21, 0x21)
GRAY_TEXT = RGBColor(0x66, 0x66, 0x66)
ACCENT_TEAL = RGBColor(0x00, 0x7B, 0x8A)         # Teal accent

HEADER_FONT = "Arial"
BODY_FONT = "Arial"

SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.50)


def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=11, font_color=WHITE, bold=False, align=PP_ALIGN.LEFT,
                     border_color=None):
    """Add a rounded rectangle with text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Pt(10)
        tf.margin_right = Pt(10)
        tf.margin_top = Pt(6)
        tf.margin_bottom = Pt(6)
        # Clear default paragraph
        p = tf.paragraphs[0]
        p.text = ""
        # Add lines
        lines = text.split("\n")
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(font_size) if i > 0 else Pt(font_size + 1)
            p.font.name = BODY_FONT
            p.font.color.rgb = font_color
            p.font.bold = bold if i == 0 else False
            p.alignment = align
            p.space_after = Pt(2)
    return shape


def add_text(slide, left, top, width, height, text, font_size=11,
             color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT):
    """Add a plain text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = BODY_FONT
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox


def add_arrow(slide, left, top, width, height):
    """Add a right-pointing arrow."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GRAY_TEXT
    shape.line.fill.background()
    return shape


def add_down_arrow(slide, left, top, width, height):
    """Add a down-pointing arrow."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GRAY_TEXT
    shape.line.fill.background()
    return shape


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # ─── Title Banner (Red) ───
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.9)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = PRIMARY_COLOR
    banner.line.fill.background()

    add_text(slide, Inches(0.5), Inches(0.1), Inches(12), Inches(0.8),
             "CI Agent Harness — Architecture & User Guide",
             font_size=28, color=WHITE, bold=True)

    # ═══════════════════════════════════════════════════════════════════
    # Section 1: HOW TO RUN IT
    # ═══════════════════════════════════════════════════════════════════
    y = 1.05
    add_text(slide, Inches(0.4), Inches(y), Inches(3), Inches(0.35),
             "1  HOW TO RUN IT", font_size=11, color=GRAY_TEXT, bold=True)

    y += 0.30
    add_rounded_rect(slide, Inches(0.4), Inches(y), Inches(4.0), Inches(0.55),
                     WHITE, "Claude Code  (recommended)",
                     font_size=12, font_color=DARK_TEXT, bold=True,
                     border_color=DARK_BLUE)

    add_rounded_rect(slide, Inches(4.8), Inches(y), Inches(3.8), Inches(0.55),
                     WHITE, "Claude App (claude.ai)",
                     font_size=12, font_color=DARK_TEXT, bold=True,
                     border_color=DARK_BLUE)

    # ═══════════════════════════════════════════════════════════════════
    # Section 2: WHAT YOU CAN ASK (5 capabilities)
    # ═══════════════════════════════════════════════════════════════════
    y += 0.75
    add_text(slide, Inches(0.4), Inches(y), Inches(6), Inches(0.35),
             "2  WHAT YOU CAN ASK  (5 core capabilities)", font_size=11, color=GRAY_TEXT, bold=True)

    y += 0.35
    box_w = Inches(2.35)
    box_h = Inches(1.15)
    gap = Inches(0.18)
    x_start = Inches(0.4)

    capabilities = [
        ("Extract Data", "Press release, CTgov,\nPubMed, pasted text"),
        ("Store in DB", "QC'd by user before\nsaving · queryable API"),
        ("Interpret BNMA", "Forest, ridge, league\ntable · auto from figures/"),
        ("Research\nLandscape", "Approved + pipeline\ncompetitors compared"),
        ("Generate Slides", "Quick (5) or Detailed\n(8+) · figures embedded"),
    ]

    for i, (title, desc) in enumerate(capabilities):
        x = x_start + i * (box_w + gap)
        shape = add_rounded_rect(slide, x, Inches(y), box_w, box_h,
                                 DARK_BLUE, f"{title}\n{desc}",
                                 font_size=9, font_color=WHITE, bold=True,
                                 align=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════════════════════════════
    # Section 3: DATA SOURCE
    # ═══════════════════════════════════════════════════════════════════
    y += 1.35
    add_text(slide, Inches(0.4), Inches(y), Inches(6), Inches(0.35),
             "3  WHERE DATA COMES FROM", font_size=11, color=GRAY_TEXT, bold=True)

    y += 0.32
    add_rounded_rect(slide, Inches(0.4), Inches(y), Inches(5.6), Inches(0.7),
                     WHITE, "Option 1: User supplies source  (faster)\nPress release URL, pasted text, or PDF pages as images",
                     font_size=10, font_color=DARK_TEXT, bold=True,
                     border_color=PRIMARY_COLOR)

    add_rounded_rect(slide, Inches(6.3), Inches(y), Inches(5.6), Inches(0.7),
                     WHITE, "Option 2: Agent searches  (broader)\nClinicalTrials.gov, press releases, PubMed, competitor DB",
                     font_size=10, font_color=DARK_TEXT, bold=True,
                     border_color=PRIMARY_COLOR)

    # ═══════════════════════════════════════════════════════════════════
    # Section 4: WORKFLOW
    # ═══════════════════════════════════════════════════════════════════
    y += 0.92
    add_text(slide, Inches(0.4), Inches(y), Inches(8), Inches(0.35),
             "4  AGENT WORKFLOW  (for slide deck generation)", font_size=11, color=GRAY_TEXT, bold=True)

    y += 0.35
    flow_h = Inches(0.85)
    flow_boxes = [
        ("Figure Inventory", "Auto-scan figures/\nClassify & route"),
        ("Extract & QC", "Structure arms, N,\nendpoints · user confirms"),
        ("Interpret & Chart", "BNMA read-out,\nlandscape comparison"),
        ("Assemble Deck", "Embed plots + figures,\nspeaker notes, QA check"),
    ]

    flow_w = Inches(2.65)
    flow_gap = Inches(0.15)
    arrow_w = Inches(0.3)
    x = Inches(0.4)

    for i, (title, desc) in enumerate(flow_boxes):
        add_rounded_rect(slide, x, Inches(y), flow_w, flow_h,
                         MID_BLUE, f"{title}\n{desc}",
                         font_size=9.5, font_color=WHITE, bold=True,
                         align=PP_ALIGN.CENTER)
        x += flow_w + flow_gap
        if i < len(flow_boxes) - 1:
            add_arrow(slide, x - Inches(0.05), Inches(y) + flow_h / 2 - Inches(0.1),
                      arrow_w, Inches(0.2))
            x += arrow_w + flow_gap

    # ═══════════════════════════════════════════════════════════════════
    # Section 5: OUTPUT + INDICATIONS
    # ═══════════════════════════════════════════════════════════════════
    y += 1.10
    add_text(slide, Inches(0.4), Inches(y), Inches(4), Inches(0.35),
             "5  OUTPUT", font_size=11, color=GRAY_TEXT, bold=True)

    y += 0.30
    add_rounded_rect(slide, Inches(0.4), Inches(y), Inches(3.8), Inches(0.65),
                     DARK_BLUE, "Quick · 5 slides\nLeadership briefing: headline,\ndata, landscape, implications",
                     font_size=9.5, font_color=WHITE, bold=True,
                     align=PP_ALIGN.CENTER)

    add_rounded_rect(slide, Inches(4.5), Inches(y), Inches(3.8), Inches(0.65),
                     DARK_BLUE, "Detailed · 8+ slides\nPresenter prep: full data, BNMA,\nstudy design, figures, backup",
                     font_size=9.5, font_color=WHITE, bold=True,
                     align=PP_ALIGN.CENTER)

    # Indications on the right side
    add_text(slide, Inches(8.8), Inches(y - 0.30), Inches(4.2), Inches(0.35),
             "SUPPORTED INDICATIONS (12)", font_size=11, color=GRAY_TEXT, bold=True)

    add_rounded_rect(slide, Inches(8.8), Inches(y), Inches(4.1), Inches(0.65),
                     LIGHT_BG,
                     "AD · Psoriasis · UC · RA · CRSwNP · PsA\nCrohn's · SLE · Asthma · COPD · IPF · AR",
                     font_size=9.5, font_color=DARK_TEXT, bold=False,
                     align=PP_ALIGN.CENTER, border_color=GRAY_TEXT)

    # ─── Footer ───
    add_text(slide, Inches(0.4), Inches(7.1), Inches(12), Inches(0.3),
             "Lilly Immunology Competitive Intelligence · Built with Claude Code · Review is required before disclosure.",
             font_size=8, color=GRAY_TEXT)

    # ─── Save ───
    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    output_path = os.path.join(base_dir, "CI_Agent_Architecture_2026-07-29.pptx")
    prs.save(output_path)
    print(f"✓ Saved: {output_path}")
    print(f"  Slides: 1")


if __name__ == "__main__":
    main()
