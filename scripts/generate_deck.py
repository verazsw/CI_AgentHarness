#!/usr/bin/env python3
"""
Competitor Landscape Slide Deck Generator
Uses python-pptx to create polished Lilly-branded presentations.

Design: White background, plain text boxes (no card panels, no footer bar).
Figures folder: All images (BNMA plots, press release pages) go in figures/

Generated: 2026-07-29
Target: Envudeucitinib (ESK-001) & Zasocitinib (TAK-279) — Phase 3 Psoriasis
Mode: Detailed presenter deck (combined)
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ═══════════════════════════════════════════════════════════════════
# PATHS
# ═══════════════════════════════════════════════════════════════════

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIGURES_DIR = os.path.join(BASE_DIR, "figures")

# ═══════════════════════════════════════════════════════════════════
# DESIGN SYSTEM CONSTANTS (13.33 × 7.50 widescreen)
# ═══════════════════════════════════════════════════════════════════

PRIMARY_COLOR = RGBColor(0xE1, 0x25, 0x1B)       # Lilly Red 2024
DARK_TEXT = RGBColor(0x21, 0x21, 0x21)            # Dark gray body text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_GOLD = RGBColor(0xFF, 0xC7, 0x09)          # Gold for highlights
SECTION_HEADER_COLOR = RGBColor(0xC8, 0x10, 0x2E) # Lilly Red for section headers
CAVEAT_COLOR = RGBColor(0x66, 0x66, 0x66)         # Medium gray for footnotes

HEADER_FONT = "Arial"
BODY_FONT = "Arial"

# Slide dimensions (widescreen 16:9)
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.50)


# ═══════════════════════════════════════════════════════════════════
# HELPER FUNCTIONS
# ═══════════════════════════════════════════════════════════════════

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=11,
                 font_name=None, color=None, bold=False, italic=False,
                 align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    font_name = font_name or BODY_FONT
    color = color or DARK_TEXT
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = font_name
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = align
    return txBox


def add_section_bullets(slide, left, top, width, height, header, bullets,
                        header_size=16, bullet_size=13, header_color=None,
                        bullet_color=None):
    """Add a section header followed by bullet points — plain text, no card."""
    header_color = header_color or SECTION_HEADER_COLOR
    bullet_color = bullet_color or DARK_TEXT
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Header
    p = tf.paragraphs[0]
    p.text = header
    p.font.size = Pt(header_size)
    p.font.name = HEADER_FONT
    p.font.color.rgb = header_color
    p.font.bold = True
    p.space_after = Pt(6)

    # Bullets
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(bullet_size)
        p.font.name = BODY_FONT
        p.font.color.rgb = bullet_color
        p.level = 0
        p.space_after = Pt(4)


def add_top_bar(slide, title):
    """Simple top bar with primary color and white title text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.0)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_COLOR
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.05), Inches(12), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.name = HEADER_FONT
    p.font.color.rgb = WHITE
    p.font.bold = True


def add_disclaimer(slide, text=None):
    """Small disclaimer text at the bottom — no bar, just text."""
    text = text or "Review is required before disclosure."
    add_text_box(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
                 text, font_size=8.5, color=CAVEAT_COLOR, italic=True)


def add_speaker_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text


def add_image_safe(slide, path, left, top, width=None, height=None):
    """Add image, print warning if missing."""
    if not os.path.exists(path):
        print(f"  Warning: Image not found: {path}", file=sys.stderr)
        add_text_box(slide, left, top, Inches(6), Inches(1),
                     f"[Image not found: {os.path.basename(path)}]",
                     font_size=12, color=CAVEAT_COLOR)
        return None
    kwargs = {"left": left, "top": top}
    if width:
        kwargs["width"] = width
    if height:
        kwargs["height"] = height
    return slide.shapes.add_picture(path, **kwargs)


# ═══════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════

def build_title_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY_COLOR)
    add_text_box(slide, Inches(1.3), Inches(1.8), Inches(10.7), Inches(2.0),
                 data["title"], font_size=36, font_name=HEADER_FONT,
                 color=WHITE, bold=True, align=PP_ALIGN.CENTER,
                 valign=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, Inches(1.3), Inches(3.9), Inches(10.7), Inches(1.0),
                 data["subtitle"], font_size=18, font_name=BODY_FONT,
                 color=RGBColor(0xF5, 0xD6, 0xB0), align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.3), Inches(5.2), Inches(10.7), Inches(1.0),
                 data["date"] + "\nCONFIDENTIAL — For Internal Use Only",
                 font_size=12, font_name=BODY_FONT, color=WHITE,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.3), Inches(6.5), Inches(10.7), Inches(0.5),
                 "Review is required before disclosure.",
                 font_size=9, font_name=BODY_FONT,
                 color=RGBColor(0xF5, 0xD6, 0xB0),
                 italic=True, align=PP_ALIGN.CENTER)
    add_speaker_notes(slide, data.get("speakerNotes", ""))


def build_content_slide(prs, data):
    """Content slide with multiple sections — white bg, plain text.
    Optionally embeds an image below the sections if 'imagePath' is provided."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, data["title"])

    for section in data["sections"]:
        x = section.get("x", 0.5)
        y = section.get("y", 1.2)
        w = section.get("w", 12.3)
        h = section.get("h", 2.0)
        add_section_bullets(slide, Inches(x), Inches(y), Inches(w), Inches(h),
                            section["header"], section["bullets"],
                            header_size=section.get("headerSize", 16),
                            bullet_size=section.get("bulletSize", 13))

    # Optional embedded image (e.g., study design figure below text)
    if "imagePath" in data and data["imagePath"]:
        pos = data.get("imagePosition", {"x": 0.5, "y": 3.5, "w": 12.3, "h": 3.5})
        add_image_safe(slide, data["imagePath"],
                       Inches(pos["x"]), Inches(pos["y"]),
                       width=Inches(pos["w"]), height=Inches(pos["h"]))

    add_disclaimer(slide, data.get("disclaimer"))
    add_speaker_notes(slide, data.get("speakerNotes", ""))


def build_image_slide(prs, data):
    """Embed an image (press release figure, study design, etc.)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, data["title"])
    image_path = data.get("imagePath", "")
    add_image_safe(slide, image_path,
                   Inches(0.4), Inches(1.15),
                   width=Inches(12.5), height=Inches(5.6))
    add_disclaimer(slide, data.get("disclaimer"))
    add_speaker_notes(slide, data.get("speakerNotes", ""))


def build_bnma_slide(prs, data):
    """Key findings on left + BNMA plot on right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, data["title"])

    # Interpretation bullets on left
    add_section_bullets(slide, Inches(0.5), Inches(1.2), Inches(4.5), Inches(5.5),
                        "Key Findings", data.get("interpretation", []),
                        header_size=16, bullet_size=13)

    # BNMA image on right
    image_path = data.get("imagePath", "")
    add_image_safe(slide, image_path,
                   Inches(5.2), Inches(1.2),
                   width=Inches(7.8), height=Inches(5.5))

    add_disclaimer(slide, "Random-effects BNMA; indirect comparison only. Review is required before disclosure.")
    add_speaker_notes(slide, data.get("speakerNotes", ""))


def build_two_column_slide(prs, data):
    """Two-column layout — left text + right text OR right image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, data["title"])

    # Left column (always text)
    left = data.get("leftColumn", {})
    if left:
        add_section_bullets(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(5.5),
                            left["header"], left["bullets"],
                            header_size=16, bullet_size=13)

    # Right side: either image or text column
    if "rightImage" in data and data["rightImage"]:
        add_image_safe(slide, data["rightImage"],
                       Inches(6.5), Inches(1.2),
                       width=Inches(6.5), height=Inches(5.5))
    elif "rightColumn" in data:
        right = data["rightColumn"]
        add_section_bullets(slide, Inches(6.8), Inches(1.2), Inches(5.8), Inches(5.5),
                            right["header"], right["bullets"],
                            header_size=16, bullet_size=13)

    add_disclaimer(slide, data.get("disclaimer"))
    add_speaker_notes(slide, data.get("speakerNotes", ""))


def build_summary_slide(prs, data):
    """Key takeaways — white background with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, data["title"] if "title" in data else "Key Takeaways")

    takeaways = data.get("takeaways", [])
    takeaway_text = "\n".join(f"• {t}" for t in takeaways)
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(3.2),
                 takeaway_text, font_size=13, font_name=BODY_FONT, color=DARK_TEXT)

    actions = data.get("actions", [])
    if actions:
        add_text_box(slide, Inches(0.8), Inches(4.7), Inches(11.5), Inches(0.5),
                     "Recommended Actions", font_size=16, font_name=HEADER_FONT,
                     color=SECTION_HEADER_COLOR, bold=True)
        actions_text = "\n".join(f"• {a}" for a in actions)
        add_text_box(slide, Inches(0.8), Inches(5.3), Inches(11.5), Inches(1.8),
                     actions_text, font_size=13, font_name=BODY_FONT, color=DARK_TEXT)

    add_speaker_notes(slide, data.get("speakerNotes", ""))


def build_table_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, data["title"])
    table_data = data.get("table", {})
    headers = table_data.get("headers", [])
    rows = table_data.get("rows", [])
    if headers:
        num_rows = len(rows) + 1
        num_cols = len(headers)
        tbl = slide.shapes.add_table(
            num_rows, num_cols,
            Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.5)
        ).table
        for i, h in enumerate(headers):
            cell = tbl.cell(0, i)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = PRIMARY_COLOR
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.name = BODY_FONT
            p.font.color.rgb = WHITE
            p.font.bold = True
        for r_idx, row in enumerate(rows):
            for c_idx, val in enumerate(row):
                cell = tbl.cell(r_idx + 1, c_idx)
                cell.text = str(val)
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(9.5)
                p.font.name = BODY_FONT
                p.font.color.rgb = DARK_TEXT
                if r_idx % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    add_disclaimer(slide, data.get("disclaimer"))
    add_speaker_notes(slide, data.get("speakerNotes", ""))


# ═══════════════════════════════════════════════════════════════════
# DATA — Envudeucitinib & Zasocitinib | Phase 3 Psoriasis | Detailed
# ═══════════════════════════════════════════════════════════════════

DATA = {
    "mode": "detailed",
    "outputFile": "Envudeucitinib_Zasocitinib_PSO_Ph3_Detailed_2026-07-29.pptx",
    "slides": [
        # ─── Slide 1: Title ───
        {
            "type": "title",
            "title": "Competitor Landscape Update:\nEnvudeucitinib & Zasocitinib",
            "subtitle": "Psoriasis — Phase 3 Oral TYK2 Inhibitor Readouts",
            "date": "July 2026",
            "speakerNotes": "This deck compares two next-generation oral selective TYK2 inhibitors that have reported Phase 3 pivotal data in moderate-to-severe plaque psoriasis. Both represent the next wave of oral therapies positioned to challenge deucravacitinib (Sotyktu) and compete with injectable biologics on efficacy.\n\nEnvudeucitinib (ESK-001, Alumis): 40 mg BID, allosteric TYK2i — ONWARD1/2 readout Jan 2026, AAD Mar 2026.\nZasocitinib (TAK-279, Takeda): 30 mg QD, allosteric TYK2i — LATITUDE 3001/3002 readout Dec 2025, AAD Mar 2026.\n\nBoth targeting NDA submission in 2026–2027. Potential market entry ~2027–2028."
        },

        # ─── Slide 2: Drug Overview — Envudeucitinib ───
        {
            "type": "content",
            "title": "Drug Overview: Envudeucitinib (ESK-001)",
            "sections": [
                {
                    "header": "Drug Profile",
                    "bullets": [
                        "Generic: envudeucitinib | Code: ESK-001",
                        "Sponsor: Alumis Inc. (Nasdaq: ALMS), South San Francisco",
                        "Modality: Oral allosteric TYK2 inhibitor (highly selective)",
                        "Dose: 40 mg BID (tablet, no fasting requirement)",
                        "Selectivity: >1-million-fold for TYK2 over JAK1/2/3",
                        "Phase: Phase 3 complete (ONWARD1 & ONWARD2); NDA planned Q4 2026",
                        "Also in development: Psoriatic arthritis, vitiligo",
                        "52-week Phase 2 OLE (STRIDE): PASI-75 78%, PASI-90 61%, PASI-100 39%"
                    ],
                    "x": 0.5, "y": 1.15, "w": 12.3, "h": 5.5
                }
            ],
            "speakerNotes": "Envudeucitinib was formerly known as ESK-001, developed by Alumis Inc. (which went public in 2024). It is NOT an AbbVie compound. Alumis is focused exclusively on autoimmune disease with TYK2 as the lead target.\n\nKey differentiator vs deucravacitinib (Sotyktu/BMS): envudeucitinib is positioned as substantially more potent with PASI-90 rates of 54-60% vs deucravacitinib's ~37% in the POETYK trials.\n\nThe 40 mg BID dosing is a potential weakness vs zasocitinib QD — adherence and patient preference may favor once-daily.\n\nSOURCES: Alumis press release Mar 28, 2026; ClinicalTrials.gov NCT06586112, NCT06588738; STRIDE 52-wk OLE (PMID 41061968)."
        },

        # ─── Slide 3: Drug Overview — Zasocitinib ───
        {
            "type": "content",
            "title": "Drug Overview: Zasocitinib (TAK-279)",
            "sections": [
                {
                    "header": "Drug Profile",
                    "bullets": [
                        "Generic: zasocitinib | Code: TAK-279 (formerly NDI-034858)",
                        "Sponsor: Takeda Pharmaceutical Company",
                        "Modality: Oral allosteric TYK2 inhibitor (highly selective)",
                        "Dose: 30 mg QD (once-daily tablet)",
                        "Selectivity: >1-million-fold for TYK2 over other JAK family",
                        "Phase: Phase 3 complete (LATITUDE 3001, 3002, Atlas); NDA planned FY2026",
                        "Also in development: Psoriatic arthritis, UC, Crohn's, SLE, vitiligo",
                        "Head-to-head: Beat deucravacitinib on PASI-100 (LATITUDE Atlas)"
                    ],
                    "x": 0.5, "y": 1.15, "w": 12.3, "h": 5.5
                }
            ],
            "speakerNotes": "Zasocitinib was in-licensed by Takeda from Nimbus Therapeutics (acquired the TYK2 program for $6B in 2022). Originally NDI-034858. Takeda has built a massive psoriasis program around it with 6+ Phase 3 studies.\n\nKey advantage: QD dosing (vs BID for envudeucitinib and deucravacitinib). Also has a broader indication portfolio than envudeucitinib.\n\nThe LATITUDE Atlas head-to-head study (NCT06973291) directly demonstrated superiority over deucravacitinib on PASI-100 at Week 16 — a significant commercial differentiator.\n\nSOURCES: Takeda press releases Dec 2025, Mar 2026, Jun 2026; ClinicalTrials.gov NCT06088043, NCT06108544, NCT06973291; Phase 2b (PMID 39167366)."
        },

        # ─── Slide 4: Mechanism of Action ───
        {
            "type": "content",
            "title": "Mechanism of Action: TYK2 Allosteric Inhibition",
            "sections": [
                {
                    "header": "TYK2 in Psoriasis Pathogenesis",
                    "bullets": [
                        "TYK2 is a JAK family kinase that mediates IL-12, IL-23, and Type I IFN signaling",
                        "IL-23/TYK2 axis is central to Th17 differentiation and IL-17 production",
                        "Selective TYK2 inhibition blocks psoriasis pathogenesis without JAK1/2/3 effects",
                        "Avoids JAK1/2/3-related risks: thrombosis, infections, cytopenias, lipid changes",
                        "Both envudeucitinib and zasocitinib bind the TYK2 pseudokinase (JH2) domain"
                    ],
                    "x": 0.5, "y": 1.15, "w": 12.3, "h": 2.2
                },
                {
                    "header": "Differentiation: Envudeucitinib vs Zasocitinib vs Deucravacitinib",
                    "bullets": [
                        "All three are allosteric TYK2 inhibitors binding the JH2 pseudokinase domain",
                        "Envudeucitinib (Alumis): 40 mg BID — highest potency, twice-daily dosing",
                        "Zasocitinib (Takeda): 30 mg QD — once-daily convenience, proven superior to deucravacitinib",
                        "Deucravacitinib (BMS/Sotyktu): 6 mg QD — first-in-class, approved but lower efficacy ceiling",
                        "Both next-gen TYK2i target PASI-90 >50% (vs ~37% for deucravacitinib)"
                    ],
                    "x": 0.5, "y": 3.5, "w": 12.3, "h": 2.5
                }
            ],
            "speakerNotes": "All three allosteric TYK2 inhibitors bind the same regulatory pseudokinase (JH2) domain — a different binding site from the catalytic JH1 domain targeted by traditional JAK inhibitors. This pseudokinase-binding mechanism enables >1-million-fold selectivity over JAK1/2/3.\n\nThe key clinical question: why are the next-gen TYK2i more efficacious?\n1. Potentially greater TYK2 occupancy at their respective doses\n2. Different pharmacokinetic profiles (longer target engagement)\n3. Trial design differences (different patient populations, estimands)\n\nDeucravacitinib POETYK data: PASI-75 ~53%, PASI-90 ~37% at Wk16.\nNext-gen target: PASI-90 >50%, approaching biologic efficacy."
        },

        # ─── Slide 5: Study Designs ───
        {
            "type": "twoColumn",
            "title": "Phase 3 Study Designs — Head-to-Head Comparison",
            "leftColumn": {
                "header": "Envudeucitinib: ONWARD1 & ONWARD2",
                "bullets": [
                    "ONWARD1: N=912 | ONWARD2: N=862",
                    "Randomized 2:1:1 (drug:PBO:apremilast)",
                    "Moderate-to-severe plaque psoriasis",
                    "24-week double-blind period",
                    "Co-primary: PASI-75 + sPGA 0/1 at Wk16 vs PBO",
                    "Key secondary: vs apremilast at Wk24",
                    "Global, multicenter, 2024–2025",
                    "ONWARD3 (OLE): N=1,680 ongoing"
                ]
            },
            "rightColumn": {
                "header": "Zasocitinib: LATITUDE 3001 & 3002",
                "bullets": [
                    "3001: N=693 | 3002: N=1,108",
                    "Randomized vs PBO vs apremilast",
                    "Moderate-to-severe plaque psoriasis",
                    "3001: 52-week | 3002: 60-week (w/ withdrawal)",
                    "Co-primary: PASI-75 + sPGA 0/1 at Wk16 vs PBO",
                    "Key secondary: vs apremilast at Wk16/24",
                    "21 countries, 2023–2025",
                    "Atlas (NCT06973291): H2H vs deucravacitinib"
                ]
            },
            "disclaimer": "Review is required before disclosure.",
            "speakerNotes": "TRIAL DESIGN COMPARISON:\n\nSimilarities:\n- Both used active comparator (apremilast) + placebo\n- Both had PASI-75 + sPGA 0/1 as co-primaries at Week 16\n- Both in moderate-to-severe plaque psoriasis (PASI≥12, BSA≥10%)\n\nDifferences:\n- Envudeucitinib: 2:1:1 randomization; shorter (24-wk double-blind)\n- Zasocitinib: longer studies (52–60 wk) with withdrawal/retreatment design in 3002\n- Zasocitinib has a dedicated H2H vs deucravacitinib (LATITUDE Atlas, N=606)\n- Envudeucitinib total program N=~1,774; Zasocitinib total N=~1,801 (PsO only)\n\nBoth programs are large, well-powered, and designed for regulatory submission."
        },

        # ─── Slide 6: Key Efficacy — Envudeucitinib ───
        {
            "type": "content",
            "title": "Key Efficacy: Envudeucitinib Phase 3 Results",
            "sections": [
                {
                    "header": "ONWARD1 (N=912) & ONWARD2 (N=862) — Week 16 & 24",
                    "bullets": [
                        "PASI-75 at Wk16: co-primary MET (p<0.001 vs PBO) — exact % not publicly disclosed",
                        "PASI-90 at Wk16: 59.9% (ONWARD1) / 53.1% (ONWARD2) vs 4.8%/4.3% PBO",
                        "PASI-100 at Wk16: 29.4% (ONWARD1) / 27.7% (ONWARD2) vs 0.9%/0.9% PBO",
                        "PASI-90 at Wk24: 68.0% (ONWARD1) / 62.1% (ONWARD2)",
                        "PASI-100 at Wk24: 41.0% (ONWARD1) / 39.5% (ONWARD2)",
                        "All ranked primary and secondary endpoints met vs PBO and apremilast"
                    ],
                    "x": 0.5, "y": 1.15, "w": 12.3, "h": 2.6
                },
                {
                    "header": "Additional Outcomes",
                    "bullets": [
                        "Scalp (ss-PGA 0/1): ~75% by Week 24 (baseline ss-PGA ≥3)",
                        "DLQI 0/1: ~50% by Week 12 (baseline DLQI ≥2)",
                        "Pruritus NRS: >4-point improvement by Wk16; clinically meaningful itch relief by Wk2",
                        "Rapid onset: >30% achieving ss-PGA 0/1 by Week 4",
                        "Sustained improvement: responses continued to increase through Wk24"
                    ],
                    "x": 0.5, "y": 3.9, "w": 12.3, "h": 2.2
                }
            ],
            "disclaimer": "Source: Alumis press release March 28, 2026. Review is required before disclosure.",
            "speakerNotes": "ENVUDEUCITINIB EFFICACY DETAIL:\n\nThe press release highlighted PASI-90 and PASI-100 rather than PASI-75 — this is a strategic choice to position above deucravacitinib and closer to biologics. The fact that PASI-90 was ~54-60% and PASI-100 ~28-29% at Wk16 suggests PASI-75 is likely >70% (must be higher than PASI-90 by definition).\n\nKey comparison to deucravacitinib (POETYK):\n- Deucravacitinib PASI-90 Wk16: ~37%\n- Envudeucitinib PASI-90 Wk16: ~54-60% → substantially higher\n\n52-week STRIDE OLE data confirms durability: PASI-75 78%, PASI-90 61%, PASI-100 39% — consistent with Phase 3 Week 24 data.\n\nSAFETY: Headache, nasopharyngitis, URTI, acne most common. No TB reactivation. No new signals."
        },

        # ─── Slide 7: Key Efficacy — Zasocitinib ───
        {
            "type": "content",
            "title": "Key Efficacy: Zasocitinib Phase 3 Results",
            "sections": [
                {
                    "header": "LATITUDE 3001 (N=693) & 3002 (N=1,108) — Week 16",
                    "bullets": [
                        "PASI-75 at Wk16: co-primary MET (p<0.001 vs PBO) — exact % not publicly disclosed",
                        "sPGA 0/1 at Wk16: 71.4% (3001) / 69.2% (3002) vs 10.7%/12.6% PBO",
                        "PASI-90 at Wk16: 61.3% (3001) / 51.9% (3002) vs 5.0%/4.0% PBO",
                        "PASI-100 at Wk16: 33.4% (3001) / 25.2% (3002) vs 0.7%/1.1% PBO",
                        "sPGA 0 at Wk16: 39.9% (3001) / 33.7% (3002) vs 0.7%/1.4% PBO",
                        "All 44 ranked secondary endpoints met (vs PBO and apremilast at Wk16/24)"
                    ],
                    "x": 0.5, "y": 1.15, "w": 12.3, "h": 2.6
                },
                {
                    "header": "LATITUDE Atlas — Head-to-Head vs Deucravacitinib (N=606)",
                    "bullets": [
                        "Primary: PASI-100 at Wk16 — zasocitinib >35% vs deucravacitinib ~14%",
                        "Zasocitinib achieved ~2.5× higher PASI-100 rates than deucravacitinib",
                        "All secondary endpoints met (PASI-90, sPGA 0) — statistically superior",
                        "Separation visible by Week 8",
                        "Durability (Wk40–60): >90% of Wk40 PASI-75/90/sPGA 0/1 responders maintained at Wk60"
                    ],
                    "x": 0.5, "y": 3.9, "w": 12.3, "h": 2.2
                }
            ],
            "disclaimer": "Source: Takeda press releases Dec 2025, Mar 2026, Jun 2026. Review is required before disclosure.",
            "speakerNotes": "ZASOCITINIB EFFICACY DETAIL:\n\nKey differentiator: the LATITUDE Atlas H2H study (NCT06973291, N=606) directly proved superiority over deucravacitinib (the only approved oral TYK2i). This is a powerful regulatory AND commercial message.\n\nZasocitinib sPGA 0/1 (~70%) is biologic-class efficacy. The PASI-90 range (52-61%) positions it similarly to envudeucitinib.\n\nApremilast comparator data (from 3001/3002):\n- sPGA 0/1: 32.1%/29.7% vs zasocitinib 71.4%/69.2% → massive advantage\n- PASI-90: 16.8%/15.9% vs zasocitinib 61.3%/51.9%\n- PASI-100: 2.9%/4.3% vs zasocitinib 33.4%/25.2%\n\nBody site data: Scalp sPGA 0/1 ~74-77%, Palmoplantar 69-71%, Nails significant.\n\nSAFETY: TEAEs 62.1%; URTI 10.1%, acne 6.5%, nasopharyngitis 6.2%. No new signals."
        },

        # ─── Slide 8: Head-to-Head Comparison ───
        {
            "type": "twoColumn",
            "title": "Envudeucitinib vs Zasocitinib — Direct Comparison",
            "leftColumn": {
                "header": "Envudeucitinib (Alumis)",
                "bullets": [
                    "Dose: 40 mg BID (twice daily)",
                    "PASI-90 Wk16: 54–60%",
                    "PASI-100 Wk16: 28–29%",
                    "PASI-100 Wk24: 40–41%",
                    "No H2H vs deucravacitinib",
                    "Smaller company (Alumis, ALMS)",
                    "NDA: Q4 2026",
                    "Pill burden: 2 tablets/day"
                ]
            },
            "rightColumn": {
                "header": "Zasocitinib (Takeda)",
                "bullets": [
                    "Dose: 30 mg QD (once daily)",
                    "PASI-90 Wk16: 52–61%",
                    "PASI-100 Wk16: 25–33%",
                    "PASI-100 Wk24: responses increasing",
                    "H2H: 2.5× deucravacitinib on PASI-100",
                    "Large pharma (Takeda, global footprint)",
                    "NDA: FY2026 (ending Mar 2027)",
                    "Pill burden: 1 tablet/day"
                ]
            },
            "disclaimer": "Cross-trial comparison only. No head-to-head data between envudeucitinib and zasocitinib. Review is required before disclosure.",
            "speakerNotes": "HEAD-TO-HEAD ASSESSMENT (cross-trial only — no direct comparison exists):\n\nEFFICACY: Very similar overall. Both achieve PASI-90 ~52-61% at Wk16 — a step-change above deucravacitinib (~37%). PASI-100 also comparable (~25-33% at Wk16).\n\nDOSING: Zasocitinib has QD advantage over envudeucitinib's BID. Patient preference and adherence favor QD, especially long-term.\n\nCOMMERCIAL: Takeda has global infrastructure, dermatology experience, and a larger program (6+ studies). Alumis is a smaller biotech that would need to build or partner for commercialization.\n\nH2H DATA: Zasocitinib has proven superiority over deucravacitinib (LATITUDE Atlas). Envudeucitinib does not have a H2H — this is a messaging gap.\n\nBOTH are clear threats to: deucravacitinib (efficacy gap), apremilast (efficacy gap), and biologics (convenience of oral)."
        },

        # ─── Slide 9: Oral Psoriasis Landscape Table ───
        {
            "type": "table",
            "title": "Psoriasis Oral Therapy Landscape — PASI Outcomes at Week 16",
            "table": {
                "headers": ["Drug", "Sponsor", "MOA", "Dose", "PASI-90", "PASI-100", "Phase"],
                "rows": [
                    ["Zasocitinib", "Takeda", "TYK2i", "30 mg QD", "52–61%", "25–33%", "Ph3 (NDA)"],
                    ["Envudeucitinib", "Alumis", "TYK2i", "40 mg BID", "53–60%", "28–29%", "Ph3 (NDA)"],
                    ["Icotrokinra*", "BMS", "TYK2i", "200 mg QD", "~55%†", "~30%†", "Phase 3"],
                    ["Deucravacitinib", "BMS", "TYK2i", "6 mg QD", "~37%", "~13%", "Approved"],
                    ["Apremilast", "Amgen", "PDE4i", "30 mg BID", "~22%", "~5%", "Approved"],
                    ["Upadacitinib‡", "AbbVie", "JAK1i", "15 mg QD", "—", "—", "Not in PsO"],
                    ["Tofacitinib‡", "Pfizer", "pan-JAK", "10 mg BID", "~40%", "~15%", "Withdrawn"]
                ]
            },
            "disclaimer": "*Icotrokinra (BMS-986165 follow-on) data estimated from early reports. †Estimated, not confirmed. ‡Non-selective JAKi included for reference only — carry BBW. Cross-trial comparisons are indirect. Review is required before disclosure.",
            "speakerNotes": "ORAL LANDSCAPE HIERARCHY (emerging from Ph3 data):\n\nTIER 1 (PASI-90 >50%): Zasocitinib, Envudeucitinib, Icotrokinra\n- These represent the next generation of oral TYK2 inhibitors\n- Approaching biologic-class efficacy in a pill\n\nTIER 2 (PASI-90 ~35-40%): Deucravacitinib\n- First-in-class, approved, but substantially lower efficacy ceiling\n- Vulnerable to next-gen displacement\n\nTIER 3 (PASI-90 <25%): Apremilast\n- PDE4 inhibitor, different MOA, good safety but limited efficacy\n- Will retain niche for milder disease / safety-concerned patients\n\nNon-selective JAK inhibitors (tofacitinib, upadacitinib) excluded from PsO market due to safety (BBW for MACE, VTE, malignancy).\n\nImplication: Both zasocitinib and envudeucitinib could significantly displace deucravacitinib market share upon approval."
        },

        # ─── Slide 10: BNMA — PASI-75 Wk16 ───
        {
            "type": "bnma",
            "title": "BNMA: PASI-75 at Week 16 (Placebo-Subtracted)",
            "imagePath": os.path.join(FIGURES_DIR, "zasocitinib_envudeucitinib_PASI75_Wk16_2026-07-29.png"),
            "interpretation": [
                "Biologics dominate: bimekizumab (~0.88), ixekizumab (~0.82), risankizumab (~0.80)",
                "Secukinumab ~0.78 — top tier of IL-17 class",
                "Icotrokinra 200mg QD ~0.63 — highest among oral TYK2i",
                "Zasocitinib Ph3 dose peaks ~0.58 (green)",
                "Envudeucitinib 40mg BID peaks ~0.55 (red/brown)",
                "Both clearly superior to deucravacitinib (~0.48) and apremilast (~0.26)",
                "Wide CrI on envudeucitinib suggests greater uncertainty"
            ],
            "speakerNotes": "BNMA INTERPRETATION — PASI-75 at Week 16:\n\nThis ridge plot shows posterior distributions of treatment difference vs placebo (random-effects BNMA).\n\nKEY FINDINGS:\n1. Injectable biologics remain superior for PASI-75: bimekizumab (0.88), ixekizumab (0.82), risankizumab (0.80), secukinumab (0.78)\n2. Among oral TYK2i: icotrokinra (BMS) leads at ~0.63, followed by zasocitinib (~0.58) and envudeucitinib (~0.55)\n3. Zasocitinib and envudeucitinib overlap substantially — no clear winner between them\n4. Both are clearly separated from deucravacitinib (~0.48) — confirms next-gen efficacy advantage\n5. Apremilast (~0.26) is far below — different efficacy tier entirely\n\nLIMITATION: Indirect comparison only. Different trials, populations, imputation methods. BNMA provides relative positioning but not definitive ranking.\n\nSOURCE: Lilly internal BNMA (Batman model), random-effects NMA, NRI estimand, mixed population."
        },

        # ─── Slide 11: BNMA — PASI-100 Wk16 ───
        {
            "type": "bnma",
            "title": "BNMA: PASI-100 at Week 16 (Placebo-Subtracted)",
            "imagePath": os.path.join(FIGURES_DIR, "zasocitinib_envudeucitinib_PASI100_Wk16_2026-07-29.png"),
            "interpretation": [
                "Bimekizumab dominates PASI-100 (~0.62) — clear separation",
                "Ixekizumab (~0.43) and risankizumab (~0.40) form second tier",
                "Secukinumab ~0.37 — below other IL-17s for complete clearance",
                "Icotrokinra ~0.33 — best oral TYK2i for PASI-100",
                "Envudeucitinib ~0.27 and zasocitinib ~0.25 — overlapping",
                "Red arrows highlight both focus compounds",
                "Deucravacitinib ~0.10 — substantially lower complete clearance"
            ],
            "speakerNotes": "BNMA INTERPRETATION — PASI-100 at Week 16:\n\nPASI-100 (complete skin clearance) is the most stringent endpoint and shows greater differentiation between agents.\n\nKEY FINDINGS:\n1. Bimekizumab (IL-17A/F) dominates at ~0.62 — best-in-class for complete clearance\n2. IL-17A inhibitors (ixekizumab ~0.43) and IL-23 (risankizumab ~0.40) form the biologic middle tier\n3. Oral TYK2i cluster: icotrokinra ~0.33, envudeucitinib ~0.27, zasocitinib ~0.25\n4. Envudeucitinib edges slightly ahead of zasocitinib for PASI-100 — but CrIs overlap\n5. Deucravacitinib (~0.10) and apremilast (~0.04) are significantly lower\n\nCLINICAL IMPLICATION: For patients prioritizing complete clearance (PASI-100), biologics remain superior to oral options. However, next-gen TYK2i are approaching secukinumab-level complete clearance rates.\n\nSOURCE: Lilly internal BNMA, random-effects, NRI estimand, mixed population."
        },

        # ─── Slide 12: BNMA — PASI-75 Wk24 ───
        {
            "type": "bnma",
            "title": "BNMA: PASI-75 at Week 24 (Placebo-Subtracted)",
            "imagePath": os.path.join(FIGURES_DIR, "zasocitinib_envudeucitinib_PASI75_Wk24_2026-07-29.png"),
            "interpretation": [
                "Risankizumab rises to #1 at Wk24 (~0.92) — sustained IL-23 effect",
                "Bimekizumab (~0.88) and ixekizumab (~0.80) remain top tier",
                "Secukinumab (~0.78), icotrokinra (~0.68) — durable responses",
                "Zasocitinib ~0.62 at Wk24 (↑ from 0.58 at Wk16)",
                "Envudeucitinib ~0.58 at Wk24 (↑ from 0.55 at Wk16)",
                "Both show continued improvement — consistent with Ph3 data",
                "Deucravacitinib ~0.53, apremilast ~0.26 — gap maintained"
            ],
            "speakerNotes": "BNMA INTERPRETATION — PASI-75 at Week 24:\n\nThe Week 24 data confirms durability and continued improvement for both focus compounds:\n\nKEY FINDINGS:\n1. Risankizumab (IL-23) surpasses all others at Wk24 (~0.92) — the long half-life and sustained IL-23 blockade enable continued improvement\n2. Bimekizumab maintains high efficacy (~0.88) at Wk24\n3. Zasocitinib improves from ~0.58 (Wk16) to ~0.62 (Wk24) — consistent with Takeda's claim of increasing responses through Wk24\n4. Envudeucitinib improves from ~0.55 (Wk16) to ~0.58 (Wk24) — consistent with Alumis reporting continued improvement\n5. The gap between biologics and oral TYK2i narrows slightly at Wk24 as oral agents continue to improve\n\nIMPLICATION: Both drugs show a durability profile — patients who respond tend to improve further. This is important for the long-term treatment paradigm.\n\nSOURCE: Lilly internal BNMA, random-effects, NRI estimand, mixed population."
        },

        # ─── Slide 13: BNMA — PASI-100 Wk24 ───
        {
            "type": "bnma",
            "title": "BNMA: PASI-100 at Week 24 (Placebo-Subtracted)",
            "imagePath": os.path.join(FIGURES_DIR, "zasocitinib_envudeucitinib_PASI100_Wk24_2026-07-29.png"),
            "interpretation": [
                "Bimekizumab sustains PASI-100 leadership (~0.60)",
                "Risankizumab (~0.50) and ixekizumab (~0.47) strengthen at Wk24",
                "Secukinumab ~0.40 — stable from Wk16",
                "Envudeucitinib rises to ~0.35 at Wk24 (from ~0.27 at Wk16)",
                "Icotrokinra ~0.34 — similar tier to envudeucitinib",
                "Zasocitinib ~0.33 (from ~0.25 at Wk16) — catches up",
                "Both drugs show meaningful Wk16→Wk24 PASI-100 improvement"
            ],
            "speakerNotes": "BNMA INTERPRETATION — PASI-100 at Week 24:\n\nComplete clearance continues to improve for both focus compounds:\n\nKEY FINDINGS:\n1. Bimekizumab remains the clear leader for complete clearance at Wk24 (~0.60)\n2. Risankizumab and ixekizumab both strengthen their Wk24 positions vs Wk16\n3. Envudeucitinib shows meaningful improvement: 0.27 → 0.35 (Wk16→Wk24), consistent with the reported 29%→41% PASI-100\n4. Zasocitinib also improves: 0.25 → 0.33, consistent with Takeda's reported increasing responses\n5. Both oral TYK2i are now competitive with icotrokinra for PASI-100 at Wk24\n\nCLINICAL CONTEXT: The Wk16→Wk24 improvement trajectory for oral TYK2i mirrors what's seen with IL-23 biologics — suggesting a gradual immunological response rather than a plateau. This is important for patient counseling (set expectations for continued improvement beyond first assessment).\n\nSOURCE: Lilly internal BNMA, random-effects, NRI estimand, mixed population."
        },

        # ─── Slide 14: Safety Comparison ───
        {
            "type": "content",
            "title": "Safety Comparison: TYK2 Inhibitor Class",
            "sections": [
                {
                    "header": "Envudeucitinib Safety (ONWARD1/2, through Wk24)",
                    "bullets": [
                        "Generally well tolerated; consistent with Phase 2 (STRIDE)",
                        "Most common AEs: headache, nasopharyngitis, URTI, acne",
                        "No clinically significant lab abnormalities",
                        "No TB reactivation cases",
                        "TEAEs mostly mild, transient, self-limited",
                        "D/C due to AE (52-wk STRIDE OLE): 3.7%"
                    ],
                    "x": 0.5, "y": 1.15, "w": 12.3, "h": 2.2
                },
                {
                    "header": "Zasocitinib Safety (LATITUDE 3001/2, through Wk16)",
                    "bullets": [
                        "TEAEs: 62.1% (zasocitinib) vs 46.9% (PBO) vs 50.5% (apremilast)",
                        "Serious TEAEs: 3.0% vs <1% PBO vs 1.5% apremilast",
                        "Most common: URTI (10.1%), acne (6.5%), nasopharyngitis (6.2%)",
                        "No new safety signals vs Phase 2b",
                        "No clinically meaningful lab parameter changes",
                        "Profile consistent across all LATITUDE studies"
                    ],
                    "x": 0.5, "y": 3.5, "w": 12.3, "h": 2.5
                }
            ],
            "disclaimer": "Cross-trial safety comparison is limited. Review is required before disclosure.",
            "speakerNotes": "SAFETY COMPARISON — TYK2i CLASS:\n\nBoth drugs show the favorable safety profile expected from selective TYK2 inhibition (no JAK1/2/3 effects):\n- No MACE/VTE/malignancy signals (unlike pan-JAK inhibitors)\n- No significant cytopenias or lipid elevations\n- No TB reactivation\n\nCOMPARISON TABLE:\n| AE | Envudeucitinib | Zasocitinib | Deucravacitinib |\n| TEAEs | Not quantified | 62.1% | ~55% (POETYK) |\n| URTI | Common | 10.1% | ~9% |\n| Acne | Common | 6.5% | ~4% |\n| Serious | Not quantified | 3.0% | ~2% |\n\nACNE: Both next-gen TYK2i show slightly higher acne rates than deucravacitinib — possibly related to greater TYK2 inhibition. Manageable and transient.\n\nNOTE: Envudeucitinib press release did not provide specific AE percentages for Ph3 — only qualitative safety summary. Zasocitinib provided more granular data.\n\nOVERALL: Both drugs have clean safety profiles appropriate for a chronic skin disease indication. No dealbreaker signals."
        },

        # ─── Slide 15: Implications for Lilly ───
        {
            "type": "content",
            "title": "Implications for Lilly Psoriasis Strategy",
            "sections": [
                {
                    "header": "Competitive Threat Assessment",
                    "bullets": [
                        "Next-gen TYK2i (zasocitinib + envudeucitinib) approach biologic efficacy in oral form",
                        "PASI-90 >50% closes the gap with ixekizumab (Taltz) for moderate psoriasis",
                        "QD oral (zasocitinib) or BID oral (envudeucitinib) vs Q2W–Q4W injections — convenience advantage",
                        "Could erode biologic initiation for new moderate-to-severe patients",
                        "Zasocitinib's H2H win vs deucravacitinib gives a strong switching narrative",
                        "Both NDA filings in 2026–2027 → market entry ~2027–2028"
                    ],
                    "x": 0.5, "y": 1.15, "w": 12.3, "h": 2.5
                },
                {
                    "header": "Mitigating Factors & Lilly Positioning",
                    "bullets": [
                        "Taltz/ixekizumab still superior on PASI-75/90/100 (biologic advantage persists)",
                        "Bimekizumab (UCB) is the primary biologic competitor, not oral TYK2i",
                        "TYK2i unlikely to displace biologics for severe/refractory patients",
                        "Lilly's biologic portfolio (Taltz) retains efficacy advantage",
                        "Oral therapies expand the treated population (new patients who refuse injections)",
                        "Monitor: whether payers position oral TYK2i as step therapy before biologics"
                    ],
                    "x": 0.5, "y": 3.8, "w": 12.3, "h": 2.5
                }
            ],
            "disclaimer": "Review is required before disclosure.",
            "speakerNotes": "LILLY STRATEGIC ASSESSMENT:\n\nTHREAT LEVEL: MODERATE for Taltz (ixekizumab).\n\nBULL CASE for oral TYK2i:\n- Patients prefer pills over injections (no injection site reactions, no cold chain)\n- Efficacy approaching biologics: PASI-90 ~55-60% is close to ixekizumab's ~70-75%\n- Payers may mandate oral TYK2i trial before biologics (step therapy)\n- Market expansion: patients who currently refuse treatment may accept an oral option\n\nBEAR CASE for oral TYK2i:\n- Biologics remain clearly superior on all PASI endpoints (especially PASI-100)\n- BID dosing (envudeucitinib) is a compliance burden for chronic disease\n- Long-term safety unknowns (still <3 years of data)\n- Payers unlikely to restrict biologics for severe patients who need maximum efficacy\n\nLILLY ACTIONS:\n1. Reinforce Taltz efficacy superiority in medical affairs communications\n2. Monitor payer step-edit policies as TYK2i launch\n3. Evaluate if Lilly should have an oral psoriasis asset (BD opportunity?)\n4. Position Taltz for severe/refractory where oral efficacy falls short"
        },

        # ─── Slide 16: Key Takeaways ───
        {
            "type": "summary",
            "title": "Key Takeaways & Recommended Actions",
            "takeaways": [
                "Two next-gen oral TYK2 inhibitors (zasocitinib QD, envudeucitinib BID) have both delivered Phase 3 PASI-90 >50% — a step-change above deucravacitinib (~37%)",
                "Zasocitinib has proven head-to-head superiority over deucravacitinib (LATITUDE Atlas: 2.5× PASI-100) — strongest competitive positioning",
                "Both remain clearly below injectable biologics on BNMA: bimekizumab > risankizumab > ixekizumab > secukinumab >> oral TYK2i",
                "NDA filings expected 2026–2027 for both; market entry likely 2027–2028 will reshape the oral psoriasis landscape",
                "Taltz (ixekizumab) retains biologic efficacy advantage but faces market-shaping risk if payers add oral TYK2i as step therapy"
            ],
            "actions": [
                "Update BNMA when full PASI-75 data disclosed (likely at/after NDA filing)",
                "Monitor FDA advisory committee meetings and approval timelines for both agents",
                "Assess payer landscape: will oral TYK2i be positioned as step therapy before biologics?",
                "Evaluate BD opportunity in oral psoriasis space for Lilly portfolio completeness"
            ],
            "speakerNotes": "SUMMARY: The psoriasis oral landscape is about to shift dramatically. Zasocitinib and envudeucitinib represent a new generation of TYK2 inhibitors that approach biologic efficacy in oral form. They will primarily threaten deucravacitinib's market share and potentially delay biologic initiation for some patients.\n\nFor Lilly (Taltz/ixekizumab):\n- Direct efficacy threat is LIMITED — biologics remain superior\n- Market dynamics threat is MODERATE — step therapy policies could delay biologic access\n- Strategic opportunity: Lilly has no oral psoriasis asset — BD/licensing could fill this gap\n\nOPEN QUESTIONS:\n1. Will regulators approve both, creating a competitive oral TYK2i market?\n2. How will payers differentiate between zasocitinib (QD) and envudeucitinib (BID)?\n3. Will long-term safety data (3-5 years) reveal any TYK2i class signals?\n4. Should Lilly evaluate an oral psoriasis asset for portfolio completeness?\n\nReview is required before disclosure."
        }
    ]
}


# ═══════════════════════════════════════════════════════════════════
# BUILD DECK
# ═══════════════════════════════════════════════════════════════════

BUILDERS = {
    "title": build_title_slide,
    "content": build_content_slide,
    "image": build_image_slide,
    "bnma": build_bnma_slide,
    "twoColumn": build_two_column_slide,
    "summary": build_summary_slide,
    "table": build_table_slide,
}


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for slide_data in DATA["slides"]:
        builder = BUILDERS.get(slide_data["type"])
        if builder:
            builder(prs, slide_data)
        else:
            print(f"Warning: Unknown slide type '{slide_data['type']}'", file=sys.stderr)

    output_file = os.path.join(BASE_DIR, DATA["outputFile"])
    prs.save(output_file)

    print(f"✓ Saved: {output_file}")
    print(f"  Mode: {DATA['mode']}")
    print(f"  Slides: {len(DATA['slides'])}")
    bnma_count = sum(1 for s in DATA["slides"] if s["type"] == "bnma")
    image_count = sum(1 for s in DATA["slides"] if s["type"] == "image")
    if image_count > 0:
        print(f"  Press release figures embedded: {image_count}")
    if bnma_count > 0:
        print(f"  BNMA plots embedded: {bnma_count}")


if __name__ == "__main__":
    main()
